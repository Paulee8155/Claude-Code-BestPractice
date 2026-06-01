#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_guide.py — Erzeugt das fusionierte Harness-Begleitdokument als Word (.docx)
und PowerPoint (.pptx) aus einer gemeinsamen, deutschen Inhaltsquelle.

Quelle: ECC-Guides (the-shortform/longform/security, ECC-WORKFLOW-GUIDE.de.md) +
BestPractice-Extras (/ecc-onboard, state/-Pattern, Karpathy). Spiegelt die VPS-weite,
globale Installation wider — ECC ist führend, RPI/context-mode entfernt.

Zahlen (Single Source of Truth):
  - Vendored ECC-Source (ecc/): 63 Agents, 249 Skills, ~80 Commands.
  - Global installiert (core-Profil): 80 Commands, 63 Agents, 21 ECC-Skills,
    rules/ecc = common + 19 Sprach-/Domain-Packs; mehr Skills via --profile full (249).

Aufruf:  python3 docs/build_guide.py
Ausgabe: docs/ECC-Harness-Guide.de.docx  und  docs/ECC-Harness-Guide.de.pptx
"""

import os

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn as pqn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
DOCX_OUT = os.path.join(HERE, "ECC-Harness-Guide.de.docx")
PPTX_OUT = os.path.join(HERE, "ECC-Harness-Guide.de.pptx")

# ---------------------------------------------------------------- Farbpalette (Word)
INDIGO   = RGBColor(0x31, 0x2E, 0x81)   # primär (Überschriften)
INDIGO_D = RGBColor(0x1E, 0x1B, 0x4B)   # tief (Deckblatt)
CYAN     = RGBColor(0x06, 0xB6, 0xD4)   # Akzent
AMBER    = RGBColor(0xB4, 0x53, 0x09)   # Hinweis/Highlight (dunkler f. Kontrast)
EMERALD  = RGBColor(0x04, 0x78, 0x57)   # OK/Erfolg
INK      = RGBColor(0x1F, 0x29, 0x37)   # Text
MUTED    = RGBColor(0x6B, 0x72, 0x80)   # gedämpft
SURFACE  = "EEF2FF"                      # heller Indigo-Hintergrund (hex str)
CODEBG   = "1E1B2E"                       # Code-Box dunkel
TBLHEAD  = "312E81"                       # Tabellenkopf
TBLZEBRA = "F3F4F6"

# ------------------------------------------------------------- Farbpalette (PPTX)
# Dark-Tech / Luxury
P_BG       = PRGB(0x0A, 0x0E, 0x1A)   # Folien-Hintergrund (fast schwarz, blaustichig)
P_CARD     = PRGB(0x14, 0x1A, 0x2E)   # Karte
P_CODE     = PRGB(0x0E, 0x14, 0x24)   # Code-Fläche / Tabellenkopf
P_CYAN     = PRGB(0x22, 0xD3, 0xEE)   # Akzent
P_AMBER    = PRGB(0xFB, 0xBF, 0x24)   # Highlight / Hinweis
P_TEXT     = PRGB(0xE6, 0xED, 0xF8)   # Haupttext hell
P_MUTED    = PRGB(0x8B, 0x98, 0xB0)   # gedämpft
P_INDIGO   = PRGB(0x81, 0x8C, 0xF8)   # sekundärer Akzent
P_EMERALD  = PRGB(0x34, 0xD3, 0x99)   # OK/Erfolg
P_WHITE    = PRGB(0xFF, 0xFF, 0xFF)

X_SHORT = "https://x.com/affaanmustafa/status/2012378465664745795"
X_LONG  = "https://x.com/affaanmustafa/status/2014040193557471352"
X_SEC   = "https://x.com/affaanmustafa/status/2033263813387223421"
X_AUTH  = "https://x.com/affaanmustafa"


# ============================================================ Word-Hilfsfunktionen
def _shade(elm, fill_hex):
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), fill_hex)
    elm.append(shd)


def shade_paragraph(p, fill_hex):
    _shade(p._p.get_or_add_pPr(), fill_hex)


def shade_cell(cell, fill_hex):
    _shade(cell._tc.get_or_add_tcPr(), fill_hex)


def set_cell_text(cell, text, bold=False, color=INK, size=9.5, fill=None, white=False):
    cell.text = ""
    p = cell.paragraphs[0]
    p.paragraph_format.space_after = Pt(2)
    p.paragraph_format.space_before = Pt(2)
    run = p.add_run(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if white else color
    if fill:
        shade_cell(cell, fill)


def add_heading(doc, text, level=1):
    """Echter Word-Heading-Style (erscheint in Navigation + TOC).
    Bei level==1 zusätzlich eine Cyan-Unterstrich-Border per pBdr."""
    p = doc.add_heading(text, level=level)
    p.paragraph_format.space_before = Pt(16 if level == 1 else 10)
    p.paragraph_format.space_after = Pt(4)
    if level == 1:
        pr = p._p.get_or_add_pPr()
        pbdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "18")
        bottom.set(qn("w:space"), "2")
        bottom.set(qn("w:color"), "06B6D4")
        pbdr.append(bottom)
        pr.append(pbdr)
    return p


def add_toc(doc):
    """Echtes, klickbares Inhaltsverzeichnis als TOC-Feld + Felder beim Öffnen aktualisieren."""
    p = doc.add_paragraph()
    run = p.add_run()
    r_elm = run._r

    fld_begin = OxmlElement("w:fldChar")
    fld_begin.set(qn("w:fldCharType"), "begin")
    r_elm.append(fld_begin)

    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = 'TOC \\o "1-2" \\h \\z \\u'
    r_elm.append(instr)

    fld_sep = OxmlElement("w:fldChar")
    fld_sep.set(qn("w:fldCharType"), "separate")
    r_elm.append(fld_sep)

    placeholder = OxmlElement("w:t")
    placeholder.text = "Inhaltsverzeichnis – in Word mit F9 aktualisieren."
    r_elm.append(placeholder)

    fld_end = OxmlElement("w:fldChar")
    fld_end.set(qn("w:fldCharType"), "end")
    r_elm.append(fld_end)

    # Felder beim Öffnen aktualisieren
    upd = OxmlElement("w:updateFields")
    upd.set(qn("w:val"), "true")
    doc.settings.element.append(upd)


def add_para(doc, text, size=10.5, color=INK, bold=False, italic=False, space=4):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(space)
    run = p.add_run(text)
    run.font.size = Pt(size)
    run.font.name = "Calibri"
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return p


def add_bullets(doc, items):
    for it in items:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.space_after = Pt(2)
        # erlaubt (label, text)-Tupel mit fettem Label
        if isinstance(it, tuple):
            r = p.add_run(it[0] + " ")
            r.font.bold = True
            r.font.size = Pt(10.5)
            r.font.color.rgb = INDIGO_D
            r2 = p.add_run(it[1])
            r2.font.size = Pt(10.5)
            r2.font.color.rgb = INK
        else:
            r = p.add_run(it)
            r.font.size = Pt(10.5)
            r.font.color.rgb = INK


def add_callout(doc, label, text, color=AMBER, fill="FEF3C7"):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.left_indent = Pt(6)
    shade_paragraph(p, fill)
    r = p.add_run(label + "  ")
    r.font.bold = True
    r.font.size = Pt(10)
    r.font.color.rgb = color
    r2 = p.add_run(text)
    r2.font.size = Pt(10)
    r2.font.color.rgb = INK


def add_code(doc, code):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(4)
    p.paragraph_format.space_after = Pt(8)
    p.paragraph_format.left_indent = Pt(4)
    shade_paragraph(p, CODEBG)
    for i, line in enumerate(code.rstrip("\n").split("\n")):
        if i > 0:
            p.add_run().add_break()
        run = p.add_run(line if line else " ")
        run.font.name = "Consolas"
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x9C, 0xE3, 0xF0)


def add_table(doc, headers, rows, widths=None):
    t = doc.add_table(rows=1, cols=len(headers))
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    t.style = "Table Grid"
    for i, h in enumerate(headers):
        set_cell_text(t.rows[0].cells[i], h, bold=True, white=True, size=9.5, fill=TBLHEAD)
    for ri, row in enumerate(rows):
        cells = t.add_row().cells
        zebra = TBLZEBRA if ri % 2 == 0 else None
        for ci, val in enumerate(row):
            set_cell_text(cells[ci], val, size=9.5, fill=zebra)
    if widths:
        for row in t.rows:
            for ci, w in enumerate(widths):
                row.cells[ci].width = Inches(w)
    return t


# ============================================================ Inhalt (gemeinsam)
MODEL_ROWS = [
    ["Exploration / Suche", "Haiku 4.5", "schnell, billig, reicht zum Finden"],
    ["Einfache, gut definierte Edits", "Sonnet 4.6", "günstige Alternative"],
    ["Standard-Coding / Multi-File", "Opus 4.8 (1M)", "Default — schnell genug via /fast"],
    ["Komplexe Architektur", "Opus 4.8 + effort high", "tiefes Reasoning"],
    ["PR-Review", "Opus / Sonnet", "versteht Kontext, fängt Nuancen"],
    ["Security-Analyse", "Opus 4.8", "darf nichts übersehen"],
    ["Doku schreiben", "Haiku / Sonnet", "einfache Struktur"],
    ["Komplexe Bugs", "Opus 4.8 + effort high", "ganzes System im Kopf halten"],
    ["Worker in Multi-Agent", "Haiku 4.5", "viele günstige Parallel-Agents"],
]

BUILDING_BLOCKS = [
    ["Skills", "Durable Kern: wiederverwendbare Workflow-Bündel, laden on-demand", "~/.claude/skills/ecc/", "langlebig"],
    ["Commands", "Slash-Einstiegspunkte (/plan, /code-review …)", "~/.claude/commands/", "Komfort-Schicht"],
    ["Subagents", "Delegierbare Spezialisten mit eigenem Tool-Scope", "~/.claude/agents/", "pro Aufgabe"],
    ["Hooks", "Trigger-Automatik an Lifecycle-Events (INAKTIV per Default)", "~/.claude/hooks/", "event-gebunden"],
    ["Rules", "„Immer befolgen\"-Leitlinien je Sprache", "~/.claude/rules/ecc/", "dauerhaft"],
    ["MCP", "Prompt-getriebene Wrapper um externe Dienste", ".mcp.json / global", "pro Session"],
    ["Instincts", "Gelernte Muster aus deinen Sessions", "homunculus/instincts/", "wächst mit dir"],
]

CMD_CHEAT = [
    ["Neues Feature planen", "/plan", "Anforderungen, Risiken, Schritte — wartet auf dein OK"],
    ["Feature umsetzen", "/feature-dev", "nutzt Skills/Agents, TDD zuerst"],
    ["Projekt ECC-ready machen", "/ecc-onboard", "Stack erkennen → 1× OK → Rules/Skills + PROJECT_RULES.md + state/"],
    ["Codebase kartieren (1×)", "/update-codemaps", "token-arme Architektur-Karten in docs/CODEMAPS/ — statt neu suchen"],
    ["Fremdes Repo verstehen", "/update-codemaps · code-tour · smart-explore", "Architektur kartieren, gefuehrte Touren, Struktur-Suche"],
    ["Code geschrieben", "/code-review  (+ /python-review, /go-review …)", "Qualität + Security der Änderungen"],
    ["Build/Tests rot", "/build-fix", "inkrementelle, minimale Fixes"],
    ["Coverage / Qualität", "/test-coverage  ·  /quality-gate", "Lücken finden, Pipeline prüfen"],
    ["Aufräumen", "/refactor-clean", "toter Code + lose Dateien raus"],
    ["Library-Doku live", "context7-MCP · docs-lookup-Agent", "aktuelle API-Doku statt Halluzination"],
    ["Session sichern/laden", "/save-session  ·  /resume-session  ·  /sessions", "über den Tag hinaus arbeiten"],
    ["Checkpoint / Nebenfrage", "/checkpoint  ·  /aside", "Stand markieren / Kontext bewahren"],
    ["Gelerntes extrahieren", "/learn-eval → /evolve → /promote", "Self-Improvement-Loop"],
    ["Modell / Kosten", "/model-route  ·  /cost-report", "billigstes ausreichendes Modell"],
    ["Harness prüfen", "/harness-audit  ·  /ecc-guide", "Config-Scorecard / Onboarding"],
    ["Parallel / Loop", "/multi-plan  ·  /loop-start  ·  /loop-status", "mehrere Modelle / Automation"],
]

# Bei DIR tatsächlich aktiv (lokal/Plugins/claude.ai-Connectoren)
MCP_ACTIVE = [
    ["PDF/Word/Excel/PPTX lesen", "markitdown", "Dokumente → Markdown (lokal, global aktiv)"],
    ["Aktuelle Library-/API-Doku", "context7", "React, Prisma, Next.js … statt veraltetem Wissen (Plugin)"],
    ["Frühere Arbeit / „schon gelöst?“", "claude-mem", "deine Memory-Quelle: /mem-search, Auto-Injektion ab 2. Session"],
    ["Kontext-Tooling", "context-mode", "lokaler MCP"],
    ["Dienste anbinden", "claude.ai-Connectoren", "n8n · Notion · Gamma · Google Drive · Microsoft Learn · Kiwi"],
    ["Bash-Token sparen", "RTK (Hook, kein MCP)", "60–90 % Ersparnis, läuft transparent global"],
]

# Empfohlen, aber NICHT installiert — bei Bedarf nachrüsten
MCP_OPTIONAL = [
    ["Web-Recherche / Discovery", "exa", "breite Suche, wenn GitHub+Docs nicht reichen"],
    ["GitHub (PRs, Issues, Code-Suche)", "github", "du nutzt die gh-CLI statt MCP"],
    ["Browser / E2E / Screenshots", "playwright", "headless Automation, visuelle Tests"],
    ["Strukturierter Wissensgraph", "memory", "Wissensgraph — NICHT nötig, claude-mem deckt das ab"],
    ["Schrittweises Reasoning", "sequential-thinking", "komplexe Mehrschritt-Probleme"],
]

SECURITY_ACTIVE = [
    ("RTK-Hook", "PreToolUse:Bash bleibt unangetastet — Token-Ersparnis global."),
    ("ECC namespace-sicher", "rules/ecc/, skills/ecc/ — keine Kollision mit Eigenem."),
    ("settings.json unberührt", "kein Auto-Edit der Startup-Config; nur additiv via ./install-vps.sh --harden."),
]
SECURITY_OPTIN = [
    ("ECC-Hooks INAKTIV", "Repo-kontrollierte Hooks nicht blind aktiv. Bewusst aktivieren nach Sichtung von hooks.json."),
    ("deny-Baseline opt-in", "./install-vps.sh --harden ergänzt rm -rf / force-push / Secret-Reads (mit Backup)."),
    ("AgentShield", "npx ecc-agentshield scan lädt externes Paket — selbst vetten, dann ausführen."),
]

# Agent-Katalog (Kap. 4) — 63 Agents insgesamt, die wichtigsten gruppiert
AGENTS_CORE = [
    ("planner", "Plan vor Code"),
    ("architect / code-architect", "System-Design"),
    ("code-explorer", "Codebase tracen"),
    ("tdd-guide", "Tests zuerst"),
    ("code-reviewer", "Qualität"),
    ("security-reviewer", "OWASP / Secrets"),
    ("build-error-resolver", "Build grün"),
    ("refactor-cleaner", "toter Code raus"),
    ("e2e-runner", "kritische Flows"),
    ("doc-updater", "Doku-Sync"),
    ("performance-optimizer", "Hot-Paths"),
    ("silent-failure-hunter", "verschluckte Fehler"),
]
AGENTS_LANG = [
    "python-reviewer", "typescript-reviewer", "react-reviewer", "go-reviewer",
    "rust-reviewer", "cpp-reviewer", "java-reviewer", "kotlin-reviewer",
    "swift-reviewer", "csharp-reviewer", "fsharp-reviewer", "django-reviewer",
    "fastapi-reviewer", "flutter-reviewer",
]
AGENTS_BUILD = [
    "go-build-resolver", "rust-build-resolver", "cpp-build-resolver",
    "java-build-resolver", "kotlin-build-resolver", "swift-build-resolver",
    "dart-build-resolver", "django-build-resolver", "react-build-resolver",
    "pytorch-build-resolver",
]
AGENTS_SPECIAL = [
    "database-reviewer", "a11y-architect", "type-design-analyzer",
    "comment-analyzer", "gan-planner / generator / evaluator", "harness-optimizer",
    "marketing-agent", "seo-specialist", "network-architect", "docs-lookup",
    "mle-reviewer",
]

# Skill-Katalog (Kap. 4) — 21 ECC-Skills (core-Profil), gruppiert
SKILLS_GROUPS = [
    ("Workflow / Qualität", "tdd-workflow · verification-loop · plankton-code-quality · production-audit · eval-harness · council · error-handling"),
    ("Codebase / Recherche", "code-tour · iterative-retrieval · skill-scout · skill-stocktake"),
    ("Testing", "e2e-testing · ai-regression-testing · windows-desktop-e2e"),
    ("Lernen / Memory", "continuous-learning(-v2) · strategic-compact · agent-introspection-debugging"),
    ("Setup / Automatik", "configure-ecc · hookify-rules · agent-sort"),
]


# ============================================================ WORD-DOKUMENT
def build_docx():
    doc = Document()
    # Basisfont
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(10.5)
    style.font.color.rgb = INK

    # Echte Heading-Styles umkonfigurieren (Navigation + TOC)
    h1 = doc.styles["Heading 1"]
    h1.font.name = "Calibri"; h1.font.size = Pt(16); h1.font.bold = True
    h1.font.color.rgb = INDIGO
    h2 = doc.styles["Heading 2"]
    h2.font.name = "Calibri"; h2.font.size = Pt(12.5); h2.font.bold = True
    h2.font.color.rgb = INDIGO_D

    sec = doc.sections[0]
    sec.left_margin = Inches(0.8)
    sec.right_margin = Inches(0.8)
    sec.top_margin = Inches(0.7)
    sec.bottom_margin = Inches(0.7)

    # ---- Deckblatt
    band = doc.add_paragraph()
    band.paragraph_format.space_before = Pt(40)
    shade_paragraph(band, "1E1B4B")
    r = band.add_run("  ENGINEERING HARNESS")
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = CYAN
    t = doc.add_paragraph()
    t.paragraph_format.space_before = Pt(10)
    rt = t.add_run("ECC × BestPractice")
    rt.font.size = Pt(34)
    rt.font.bold = True
    rt.font.color.rgb = INDIGO
    sub = doc.add_paragraph()
    rs = sub.add_run("Das fusionierte, VPS-weite Claude-Code-Harness — Bedienungsanleitung")
    rs.font.size = Pt(15)
    rs.font.color.rgb = INK
    meta = doc.add_paragraph()
    rm = meta.add_run("Global aktiv in ~/.claude · ECC 2.0.0-rc.1 (core, führend) · RTK-sicher\n"
                      "Quelle/Re-Install: /root/projekte/Claude Code BestPractice  →  ./install-vps.sh")
    rm.font.size = Pt(10)
    rm.font.color.rgb = MUTED
    add_callout(doc, "Halte mich offen.",
                "Dieses Dokument ist dein Nachschlagewerk neben dem Code: wann welcher Command, "
                "welches Modell, welcher MCP. Basiert auf den Guides von @affaanmustafa (X-Links am Ende).",
                color=INDIGO_D, fill=SURFACE)

    # ---- Inhalt (echtes TOC-Feld)
    add_heading(doc, "Inhalt", 1)
    add_toc(doc)

    # 1
    add_heading(doc, "1 · Was ist dieses Harness?", 1)
    add_para(doc, "Dieses Harness ist die Fusion zweier Systeme — mit klarem Fokus auf ECC:")
    add_bullets(doc, [
        ("ECC (Everything Claude Code)", "von Affaan Mustafa — die dominante Engine: 63 Agents, 249 Skills, "
         "~80 Commands, manifest-basierter Installer, Continuous Learning, Security-Scanning. "
         "Global installiert (core-Profil): 80 Commands, 63 Agents, 21 ECC-Skills, "
         "rules/ecc = common + 19 Sprach-/Domain-Packs; mehr Skills via --profile full (249)."),
        ("BestPractice-Extras", "die wenigen einzigartigen Stärken deines alten Harness, die ECC ergänzen: "
         "/ecc-onboard (One-Shot-Projekt-Setup), das state/-Pattern und die Karpathy-Prinzipien."),
    ])
    add_para(doc, "Philosophie in einem Satz:", bold=True, space=2)
    add_callout(doc, "„Konfiguration ist Fine-Tuning, nicht Architektur.“",
                "Baue wiederverwendbare Muster, halte das Kontextfenster sauber, delegiere an das billigste "
                "ausreichende Modell, verifiziere mit Evidenz — und lass Komfort nie die Sicherheit überholen.",
                color=INDIGO_D, fill=SURFACE)
    add_para(doc, "Fünf Kernprinzipien (SOUL.md): Agent-First · Test-Driven · Security-First · "
                  "Immutability · Plan-Before-Execute.", italic=True, color=MUTED)

    # 2
    add_heading(doc, "2 · Setup & VPS-Architektur", 1)
    add_para(doc, "ECC ist jetzt GLOBAL nach ~/.claude installiert und damit in JEDEM Projekt auf der VPS aktiv "
                  "(WMS, Jarvis, n8n, Werkstatt …). Die Installation ist bewusst additiv und namespace-sicher:")
    add_bullets(doc, [
        ("Namespace-sicher:", "ECC liegt unter rules/ecc/ und skills/ecc/ — keine Kollision mit Eigenem."),
        ("RTK bleibt König:", "Der globale PreToolUse:Bash-Hook (RTK, 60–90 % Token-Ersparnis) wurde NICHT angefasst."),
        ("Startup-Config additiv:", "Der Installer lässt settings.json/CLAUDE.md unberührt. Bewusste Aufräum-Schritte "
         "(feature-dev-Plugin deaktiviert, context-mode entfernt) wurden separat mit Backup vorgenommen."),
        ("Reproduzierbar:", "Quelle ist das Repo; Re-Install via ./install-vps.sh (idempotent, mit Backup)."),
    ])
    add_para(doc, "Koexistenz mit deinem bestehenden Setup:", bold=True, space=2)
    add_table(doc, ["Dein Tool", "ECC-Pendant", "Empfehlung"], [
        ["superpowers", "tdd-workflow, verification-loop …", "Beides nutzbar; bei Doppelung ECC bevorzugen"],
        ["claude-mem", "/save-session, continuous-learning-v2", "primäre Memory-Quelle; injiziert Kontext automatisch"],
        ["RTK", "Token-Ökonomie (§10)", "Ergänzen sich: RTK auf Bash-, ECC auf Workflow-Ebene"],
        ["codex / superpowers", "diverse Skills/Agents", "eigenständiger Mehrwert; bleiben aktiv"],
    ], widths=[1.4, 2.4, 3.0])
    add_callout(doc, "Aufgeräumt:", "feature-dev existierte 3-fach — die ECC-Version ist kanonisch, das gleichnamige "
                "Plugin wurde deaktiviert. RPI wurde entfernt (redundant zu ECC). Bei Namens-Zweifeln eindeutige "
                "ECC-Commands nutzen (/harness-audit, /ecc-guide, /ecc-onboard, /instinct-status).")

    # 3
    add_heading(doc, "3 · Mentales Modell — die Bausteine", 1)
    add_table(doc, ["Baustein", "Was es ist", "Wo (global)", "Lebensdauer"], BUILDING_BLOCKS,
              widths=[0.95, 3.2, 1.7, 1.0])
    add_callout(doc, "Wichtigste Regel:", "Die durable Logik gehört in Skills. Commands sind nur der bequeme "
                "Einstieg. Was du zum dritten Mal tippst → mach eine Skill draus (/skill-create).",
                color=INDIGO_D, fill=SURFACE)

    # 4 (NEU)
    add_heading(doc, "4 · Agent- & Skill-Katalog", 1)
    add_para(doc, "63 Agents insgesamt — die wichtigsten:", bold=True, color=INDIGO, space=4)
    add_para(doc, "Core-Workflow:", bold=True, space=2)
    add_bullets(doc, [(a, f"— {d}") for a, d in AGENTS_CORE])
    add_para(doc, "Sprach-Reviewer:", bold=True, space=2)
    add_para(doc, " · ".join(AGENTS_LANG), size=10, color=INK, space=4)
    add_para(doc, "Build-Resolver:", bold=True, space=2)
    add_para(doc, " · ".join(AGENTS_BUILD), size=10, color=INK, space=4)
    add_para(doc, "Spezial:", bold=True, space=2)
    add_para(doc, " · ".join(AGENTS_SPECIAL), size=10, color=INK, space=6)

    add_para(doc, "21 ECC-Skills (core-Profil) — gruppiert:", bold=True, color=INDIGO, space=4)
    add_table(doc, ["Gruppe", "Skills"], [[g, s] for g, s in SKILLS_GROUPS], widths=[1.9, 4.9])
    add_callout(doc, "Mehr verfügbar:", "Mehr (249 Skills) via ./install-vps.sh --profile full. "
                "claude-mem (mem-search, smart-explore) + superpowers koexistieren.",
                color=INDIGO_D, fill=SURFACE)

    # 5
    add_heading(doc, "5 · Täglicher Workflow — die 5-Phasen-Pipeline", 1)
    add_para(doc, "Der Erfinder arbeitet in klaren Phasen, mit /clear dazwischen und Zwischenergebnissen in Dateien "
                  "(statt alles im Kontext zu halten):")
    add_code(doc,
             "Phase 1  RESEARCH   →  Explore-Agent / claude-mem         →  research.md\n"
             "Phase 2  PLAN       →  /plan  (wartet auf dein OK!)        →  plan.md\n"
             "Phase 3  IMPLEMENT  →  /feature-dev + Skill tdd-workflow   →  Code + Tests\n"
             "Phase 4  REVIEW     →  /code-review (+ /<sprache>-review)  →  review.md\n"
             "Phase 5  VERIFY     →  Tests/Build (+ /build-fix)          →  grün, sonst zurück zu Phase 3")
    add_bullets(doc, [
        "Jeder Schritt: ein klarer Input, ein klarer Output (in Datei speichern).",
        "Outputs werden Inputs der nächsten Phase. Phasen nicht überspringen.",
        "/clear zwischen großen Phasen — Explorations-Kontext ist für die Umsetzung irrelevant.",
        "Subagents immer objektiven Kontext mitgeben, nicht nur die Frage.",
    ])

    # 6
    add_heading(doc, "6 · Onboarding & Codebase einmal erkunden (Codemaps)", 1)
    add_para(doc, "Die teuerste Verschwendung ist, die Codebase in jeder Session neu zu durchsuchen. ECC-Prinzip: "
                  "EINMAL erkunden, das Ergebnis in Dateien festhalten, danach immer die kompakte Karte lesen statt "
                  "den ganzen Code. Drei Ebenen — projektweit einrichten, Architektur kartieren, pro Aufgabe recherchieren:",
                  space=2)

    add_para(doc, "a) Projekt ECC-ready machen — /ecc-onboard (einmalig je Projekt)", bold=True, color=INDIGO, space=2)
    add_para(doc, "Ein Command, ein OK: erkennt den Stack, installiert ECC project-level und legt persistente "
                  "Projekt-Dateien an, die Claude künftig zuerst liest.")
    add_table(doc, ["Schritt", "Was passiert", "Ergebnis (persistent)"], [
        ["1 · Detect", "Stack aus package.json / go.mod / pyproject.toml … erkennen", "Profil + Sprach-Packs"],
        ["2 · Dry-Run", "install-apply --target claude-project --dry-run", "Plan ohne Schreiben"],
        ["3 · Bestätigung", "1× OK des Users", "Freigabe"],
        ["4 · Install", "install-apply --target claude-project", ".claude/rules/ecc/ + skills/ecc/"],
        ["5 · Kontext", "Templates mit echten Werten füllen", "PROJECT_RULES.md + state/*"],
    ], widths=[1.3, 3.4, 2.1])
    add_para(doc, "Aus dem Zielprojekt-Verzeichnis ausführen (Install landet im aktuellen cwd). Idempotent: erneut "
                  "ausführen merged, überschreibt nichts blind.", italic=True, color=MUTED, size=9.5)

    add_para(doc, "b) Architektur kartieren — /update-codemaps  (KERN-FEATURE)", bold=True, color=INDIGO, space=2)
    add_para(doc, "Scannt die Struktur EINMAL und schreibt token-arme Architektur-Karten auf Platte. In späteren "
                  "Sessions liest Claude die kompakte Karte (billig) statt die ganze Codebase neu zu durchsuchen.")
    add_table(doc, ["Datei (in docs/CODEMAPS/)", "Inhalt"], [
        ["architecture.md", "System-Diagramm, Service-Grenzen, Datenfluss"],
        ["backend.md", "Routes → Controller → Service → Repository"],
        ["frontend.md", "Page-Tree, Komponenten-Hierarchie, State-Flow"],
        ["data.md", "Tabellen, Relationen, Migrations-Historie"],
        ["dependencies.md", "externe Dienste, Third-Party-Integrationen"],
    ], widths=[2.4, 4.4])
    add_callout(doc, "Wann neu bauen?", "Codemaps sind ein Snapshot — sie aktualisieren sich nicht von selbst. Nach "
                "größeren Architektur-Änderungen /update-codemaps erneut laufen lassen. Erstes Mal in einem fremden "
                "Repo: code-tour-Skill bzw. claude-mem smart-explore/learn-codebase erzeugen Touren/Struktur.",
                color=INDIGO_D, fill=SURFACE)

    add_para(doc, "c) Pro Aufgabe recherchieren — research.md-Pattern", bold=True, color=INDIGO, space=2)
    add_bullets(doc, [
        ("Phase-1-Output in Datei:", "Ein Explore-/code-explorer-Agent kartiert die für DIE Aufgabe relevanten "
         "Dateien/Muster → research.md. Das wird Input für /plan."),
        ("code-tour (optional):", "Erzeugt wiederverwendbare .tour-Dateien (öffnen direkt Datei+Zeile) — ideal für "
         "Onboarding neuer Mitarbeiter oder Architektur-Walkthroughs."),
        ("claude-mem ergänzt:", "Injiziert relevanten Kontext aus früheren Sessions ab der 2. Session automatisch — "
         "Erinnerung ohne Zutun. /mem-search für gezielte Treffer („schon mal gelöst?\")."),
    ])
    add_callout(doc, "Die Regel dahinter:", "„In Dateien speichern, nicht im Kopf behalten.\" research.md → plan.md → "
                "review.md, dazwischen /clear. Codemaps + PROJECT_RULES.md + state/ sind die dauerhafte Landkarte; "
                "der Chat-Kontext bleibt schlank.", color=INDIGO_D, fill=SURFACE)

    # 7 (NEU · Playbook)
    add_heading(doc, "7 · Playbook — Von 0 auf Feature (durchgespielt)", 1)
    add_para(doc, "Kapitel 5–6 erklären die Phasen und das Onboarding einzeln. Hier laufen sie als EIN "
                  "durchgehender Ablauf zusammen — genau wie der Erfinder: zwei Startpunkte, ein Feature von "
                  "Anfang bis grün, dann der tägliche Rhythmus.", space=2)

    add_para(doc, "a) Zwei Startpunkte — leeres Projekt ODER bestehende Codebase", bold=True, color=INDIGO, space=2)
    add_para(doc, "Leeres Projekt (greenfield) — bei null anfangen:", bold=True, space=2)
    add_code(doc,
             "mkdir mein-projekt && cd mein-projekt && git init\n"
             "claude                       # Claude Code IN diesem Ordner starten\n"
             "/ecc-onboard                 # Stack-Frage → 1x OK → .claude/rules+skills, PROJECT_RULES.md, state/\n"
             "/plan  Baue <Feature X> …    # erster Plan — wartet auf dein OK")
    add_para(doc, "Bestehende Codebase übernehmen — fremden Code zähmen:", bold=True, space=2)
    add_code(doc,
             "git clone <repo> && cd <repo>\n"
             "claude\n"
             "/ecc-onboard                 # erkennt Stack aus package.json / go.mod / pyproject.toml …\n"
             "/update-codemaps             # EINMAL kartieren → docs/CODEMAPS/*  (danach Karte statt Volltext)\n"
             "code-tour · smart-explore    # geführte Tour / Struktur-Suche bei fremdem Code")
    add_callout(doc, "Einmal, nicht jedes Mal:", "/ecc-onboard und /update-codemaps laufen pro Projekt EINMAL. "
                "Das Ergebnis ist persistent (Dateien) — jede Folge-Session startet mit der kompakten Karte statt "
                "die ganze Codebase neu zu durchsuchen.", color=INDIGO_D, fill=SURFACE)

    add_para(doc, "b) Ein Feature von Anfang bis Ende — Beispiel: neuer /api/search-Endpoint",
             bold=True, color=INDIGO, space=2)
    add_para(doc, "Eine Aufgabe, fünf Phasen, je ein Datei-Output, /clear dazwischen, billigstes ausreichendes Modell:",
             space=2)
    add_table(doc, ["Phase", "Du tippst", "Modell", "Output (Datei)"], [
        ["1 · RESEARCH", "Explore-Agent: Routing/Handler/Tests für Suche kartieren (+ claude-mem)", "Haiku", "research.md"],
        ["— /clear —", "Explorations-Kontext ist für die Umsetzung irrelevant", "—", "—"],
        ["2 · PLAN", "/plan  /api/search: Query → Filter → Paginierung, Tests zuerst", "Opus", "plan.md (wartet auf OK)"],
        ["3 · IMPLEMENT", "/feature-dev  plan.md umsetzen (Skill tdd-workflow: Tests zuerst)", "Sonnet→Opus", "Code + Tests (RED→GREEN)"],
        ["4 · REVIEW", "/code-review  (+ /python-review | /go-review …)", "Opus", "review.md"],
        ["5 · VERIFY", "Tests/Build grün? sonst /build-fix → zurück zu Phase 3", "Sonnet", "grün → commit / PR"],
    ], widths=[1.15, 3.15, 0.95, 1.55])
    add_para(doc, "Konkret getippt (copy-paste-tauglich):", bold=True, space=2)
    add_code(doc,
             "# Phase 1 — Research (Output in Datei, nicht im Kopf)\n"
             "Explore-Agent: kartiere vorhandenes Routing/Handler/Tests für Suche → research.md\n"
             "/clear\n"
             "# Phase 2 — Plan (wartet auf dein OK)\n"
             "/plan  Neuer Endpoint /api/search: Volltext-Query, Filter, Paginierung; Tests zuerst. Quelle: research.md\n"
             "/clear\n"
             "# Phase 3 — Implement (TDD: erst RED, dann GREEN)\n"
             "/feature-dev  Setze plan.md um — Tests zuerst\n"
             "# Phase 4 — Review\n"
             "/code-review                 # + /python-review bzw. /<sprache>-review fuer die geaenderten Dateien\n"
             "# Phase 5 — Verify\n"
             "/build-fix                   # nur falls Build/Tests rot — sonst committen\n"
             "/save-session                # Stand sichern (fuer morgen)")
    add_callout(doc, "Die eine Regel:", "„In Dateien speichern, nicht im Kopf behalten.“ research.md → plan.md → "
                "review.md. Jede Phase: ein Input, ein Output. Phasen nicht überspringen, /clear dazwischen.",
                color=INDIGO_D, fill=SURFACE)

    add_para(doc, "c) Der tägliche Rhythmus", bold=True, color=INDIGO, space=2)
    add_table(doc, ["Wann", "Was du tust", "Warum"], [
        ["Session-Start", "claude-mem injiziert Kontext automatisch (ab 2. Session) · rtk gain", "nahtlos weiter, wo du warst"],
        ["Vor großer Aufgabe", "/resume-session · /mem-search „schon gelöst?“", "Alt-Wissen nutzen statt neu lösen"],
        ["Während", "in Phasen arbeiten · /clear zwischen Phasen · /checkpoint", "Kontext schlank, billiges Modell"],
        ["Nebenfrage", "/aside  (oder /fork)", "Haupt-Task-Kontext nicht verlieren"],
        ["Session-Ende", "/learn-eval → /evolve → /promote · /save-session", "ECC lernt — morgen schneller"],
    ], widths=[1.5, 3.4, 1.9])
    add_callout(doc, "Der Loop, der ECC ausmacht:", "Was du zum 3. Mal tippst → /skill-create. Bewährte Muster → "
                "/promote (global). So wirst du mit jeder Session schneller — genau das trennt ECC von „nur Configs“.",
                color=INDIGO_D, fill=SURFACE)

    # 8
    add_heading(doc, "8 · Cheat-Sheet — wann welcher Command", 1)
    add_table(doc, ["Situation", "Command", "Zweck"], CMD_CHEAT, widths=[1.9, 2.5, 2.4])
    add_para(doc, "Sprach-spezifisch: /go-build|review|test, /rust-*, /cpp-*, /kotlin-*, /react-*, /flutter-*, "
                  "/python-review, /fastapi-review.  PRP-Pipeline: /prp-prd → /prp-plan → /prp-implement → "
                  "/prp-commit → /prp-pr.", color=MUTED, size=9.5)

    # 9
    add_heading(doc, "9 · Wann welcher MCP", 1)
    add_para(doc, "Faustregel: 20–30 MCPs konfiguriert, aber < 10 aktiv / < 80 Tools. Ungenutztes pro Projekt mit "
                  "/mcp deaktivieren — das Kontextfenster ist kostbar.")
    add_para(doc, "Bei dir aktiv (global):", bold=True, color=INDIGO, space=2)
    add_table(doc, ["Aufgabe", "MCP / Tool", "Hinweis"], MCP_ACTIVE, widths=[2.4, 1.7, 2.7])
    add_para(doc, "Optional / empfohlen — NICHT installiert, bei Bedarf nachrüsten:", bold=True, color=MUTED, space=2)
    add_table(doc, ["Aufgabe", "MCP / Tool", "Hinweis"], MCP_OPTIONAL, widths=[2.4, 1.7, 2.7])
    add_callout(doc, "Memory-Klarstellung:", "Der „memory“-MCP (Wissensgraph) ist NICHT installiert und nicht nötig — "
                "claude-mem ist deine primäre Memory-Quelle (§11). EINE Quelle pro Projekt, sonst doppelte Wahrheit.",
                color=INDIGO_D, fill=SURFACE)

    # 10
    add_heading(doc, "10 · Token-Ökonomie & Modellstrategie", 1)
    add_para(doc, "Billigstes ausreichendes Modell pro Aufgabe — das ist die Subagent-Architektur:")
    add_table(doc, ["Aufgabe", "Modell", "Warum"], MODEL_ROWS, widths=[2.4, 1.6, 2.8])
    add_callout(doc, "Default = Opus 4.8 (1M)", "— schnell genug via /fast. Effort-Level: medium Standard, high bei "
                "schweren Bugs/Architektur. Hilfe: /model-route, /cost-report.",
                color=INDIGO_D, fill=SURFACE)
    add_bullets(doc, [
        "Strategisches Compacting: Auto-Compact aus, manuell /compact mit Hint an logischen Punkten.",
        "Modulare Dateien (hunderte statt tausende Zeilen) → bessere First-Try-Trefferquote.",
        "/refactor-clean nach langen Sessions.",
    ])

    # 11
    add_heading(doc, "11 · Memory & Sessions", 1)
    add_para(doc, "claude-mem ist deine primäre, persistente Memory-Quelle — so nutzt du sie:",
             bold=True, color=INDIGO, space=2)
    add_bullets(doc, [
        ("Automatisch:", "injiziert relevanten Kontext aus früheren Sessions ab der 2. Session je Projekt — ohne Zutun."),
        ("/mem-search „…“:", "gezielt frühere Arbeit finden („schon gelöst?“, „wie haben wir X gemacht?“)."),
        ("Erfasst von selbst:", "Beobachtungen/Entscheidungen werden automatisch gespeichert — kein manuelles Speichern."),
    ])
    add_para(doc, "Sessions explizit sichern/laden (ECC-Commands, ergänzend):", bold=True, color=INDIGO, space=2)
    add_table(doc, ["Command", "Zweck"], [
        ["/save-session", "Stand nach ~/.claude/session-data/ sichern"],
        ["/resume-session", "letzte Session laden & weitermachen"],
        ["/sessions", "Historie durchsuchen/verwalten"],
        ["/checkpoint", "Checkpoint im laufenden Lauf"],
        ["/aside", "Nebenfrage ohne Kontextverlust"],
    ], widths=[2.0, 4.6])
    add_para(doc, "Gute Session-Datei: Was funktioniert hat (mit Evidenz) · Was nicht ging · Was offen ist.", space=2)
    add_callout(doc, "EINE Memory-Quelle:", "claude-mem ist gesetzt — den generischen „memory“-MCP NICHT zusätzlich "
                "aktivieren (sonst doppelte Wahrheit). Keine Secrets in Memory.")

    # 12
    add_heading(doc, "12 · Continuous Learning (Self-Improvement)", 1)
    add_para(doc, "Das unterscheidet ECC von „nur Configs\": es lernt aus deinen Sessions.")
    add_table(doc, ["Command", "Zweck"], [
        ["/learn", "Muster aus der Session extrahieren"],
        ["/learn-eval", "extrahieren + Qualität selbst bewerten"],
        ["/evolve", "Instincts analysieren → bessere Skill-Struktur"],
        ["/promote", "Projekt-Instincts global befördern"],
        ["/instinct-status", "alle Instincts mit Confidence-Score"],
        ["/skill-create", "aus Git-History eine Skill generieren"],
    ], widths=[2.0, 4.6])
    add_para(doc, "Routine: Session-Ende → /learn-eval → bei guten Mustern /evolve → Bewährtes mit /promote global "
                  "verfügbar machen.", space=2)

    # 13
    add_heading(doc, "13 · Parallelisierung", 1)
    add_bullets(doc, [
        ("/fork", "Konversation forken für nicht-überlappende Nebenaufgaben (Fragen zur Codebase)."),
        ("Git-Worktrees", "für überlappende parallele Arbeit ohne Konflikte — eigene Claude-Instanz je Worktree."),
        ("Cascade-Methode", "neue Tasks in neuen Tabs, alt→neu abarbeiten, max. 3–4 gleichzeitig."),
        ("ECC-Multi-Model", "/multi-plan, /multi-workflow, /multi-backend, /multi-frontend, /loop-start."),
    ])
    add_code(doc, "git worktree add ../feature-a feature-a\ncd ../feature-a && claude   # eigene Instanz")

    # 14
    add_heading(doc, "14 · Security im Alltag + VPS-Hinweise", 1)
    add_para(doc, "Kernhaltung: Baue so, als würde das Modell irgendwann etwas Feindliches lesen, während es etwas "
                  "Wertvolles hält (lethal trifecta: private Daten + untrusted Content + externe Kommunikation).")
    add_para(doc, "Bereits aktiv / sicher:", bold=True, color=EMERALD, space=2)
    add_bullets(doc, SECURITY_ACTIVE)
    add_para(doc, "Bewusst NICHT automatisch (Least-Agency, opt-in):", bold=True, color=AMBER, space=2)
    add_bullets(doc, SECURITY_OPTIN)
    add_callout(doc, "VPS-weit denken:", "Da ECC global ist, wirkt es in allen Produktionsprojekten. ECC-Hooks "
                "bleiben deshalb inaktiv, bis du sie pro Bedarf sichtest und aktivierst. Härtung gezielt via "
                "./install-vps.sh --harden (legt vorher ein settings.json-Backup an).")

    # 15
    add_heading(doc, "15 · ECC-CLIs & Wartung", 1)
    add_code(doc,
             "cd \"/root/projekte/Claude Code BestPractice/ecc\"\n"
             "node scripts/ecc.js doctor          # Health-Check (claude-home: OK)\n"
             "node scripts/ecc.js list-installed   # was ist installiert\n"
             "node scripts/ecc.js repair           # ECC-Dateien wiederherstellen\n"
             "npx ecc consult \"security reviews\"  # Advisor: welche Komponenten?\n"
             "# Re-Install / Update VPS-weit:\n"
             "cd \"/root/projekte/Claude Code BestPractice\" && ./install-vps.sh")
    add_para(doc, "Slash-Einstieg: /ecc-guide (Onboarding), /harness-audit (Config-Scorecard), /cost-report, "
                  "/skill-health, /projects.", color=MUTED, size=9.5)

    # 16
    add_heading(doc, "16 · Quellen & TL;DR", 1)
    add_para(doc, "Original-Guides von @affaanmustafa (als X-Threads veröffentlicht; lokal im Repo unter ecc/):",
             bold=True, space=2)
    add_bullets(doc, [
        ("Shortform-Guide:", X_SHORT + "  (ecc/the-shortform-guide.md) — Setup & Philosophie"),
        ("Longform-Guide:", X_LONG + "  (ecc/the-longform-guide.md) — Token, Memory, Evals, Parallelisierung"),
        ("Security-Guide:", X_SEC + "  (ecc/the-security-guide.md) — Agent-Security, Pflichtlektüre"),
        ("Erfinder:", X_AUTH + "  (@affaanmustafa)"),
    ])
    add_para(doc, "TL;DR — so arbeitest du „wie der Erfinder“:", bold=True, color=INDIGO, space=2)
    add_bullets(doc, [
        "In Phasen denken: Research → /plan → /feature-dev+TDD → /code-review → Verify, /clear dazwischen.",
        "Kontext sauber: < 10 MCPs aktiv, billigstes ausreichendes Modell, modulare Dateien, strategisch compacten.",
        "Persistieren: /save-session am Ende, /resume-session am Anfang.",
        "Lernen lassen: /learn-eval → /evolve → /promote.",
        "Security zuerst: RTK aktiv, ECC-Hooks/AgentShield bewusst aktivieren, nie blind ausführen.",
        "Bei Unsicherheit: /ecc-guide, ecc doctor, die drei Original-Guides.",
    ])

    doc.save(DOCX_OUT)
    print("✅ Word:", DOCX_OUT)


# ============================================================ POWERPOINT (Dark-Tech)
def _bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = P_BG


def _round(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(l), PInches(t), PInches(w), PInches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = PPt(0.75)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    return sh


def _rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(l), PInches(t), PInches(w), PInches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _box(slide, l, t, w, h):
    return slide.shapes.add_textbox(PInches(l), PInches(t), PInches(w), PInches(h)).text_frame


def _set(tf, text, size=18, color=P_TEXT, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = PPt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def _bullet(tf, text, size=15, color=P_TEXT, bold=False, level=0, first=False, font="Calibri", align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = PPt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def _chrome(slide, kicker, n):
    """Cyan-Punkt + Kicker oben links, Foliennummer oben rechts (n / 18)."""
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, PInches(0.6), PInches(0.42), PInches(0.16), PInches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = P_CYAN; dot.line.fill.background(); dot.shadow.inherit = False
    tf = _box(slide, 0.85, 0.33, 8.0, 0.5)
    _set(tf, kicker.upper(), size=12, color=P_CYAN, bold=True)
    tf2 = _box(slide, 10.5, 0.33, 2.2, 0.5)
    _set(tf2, f"{n} / 18", size=12, color=P_MUTED, align=PP_ALIGN.RIGHT)


def _h(slide, title):
    """Titel + Cyan-Regel darunter."""
    tf = _box(slide, 0.6, 0.85, 12.1, 0.85)
    _set(tf, title, size=27, color=P_TEXT, bold=True)
    _rect(slide, 0.65, 1.62, 2.2, 0.045, P_CYAN)


def _dtable(slide, headers, rows, l, t, w, h, col_w=None, fsize=12):
    rcount = len(rows) + 1
    ccount = len(headers)
    gtbl = slide.shapes.add_table(rcount, ccount, PInches(l), PInches(t), PInches(w), PInches(h)).table
    gtbl.first_row = False
    gtbl.horz_banding = False
    # "No Style, No Grid" Tabellen-Style versuchen
    try:
        STYLE_ID = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
        tbl = gtbl._tbl
        tblPr = tbl.tblPr
        existing = tblPr.find(pqn("a:tableStyleId"))
        if existing is None:
            sid = etree.SubElement(tblPr, pqn("a:tableStyleId"))
            tblPr.remove(sid)
            tblPr.insert(0, sid)
        else:
            sid = existing
        sid.text = STYLE_ID
    except Exception:
        pass
    if col_w:
        for i, cw in enumerate(col_w):
            gtbl.columns[i].width = PInches(cw)
    for ci, htxt in enumerate(headers):
        c = gtbl.cell(0, ci)
        c.fill.solid(); c.fill.fore_color.rgb = P_CODE
        c.margin_top = PInches(0.04); c.margin_bottom = PInches(0.04)
        c.margin_left = PInches(0.08); c.margin_right = PInches(0.08)
        tf = c.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = htxt
        r.font.size = PPt(fsize); r.font.bold = True; r.font.color.rgb = P_CYAN; r.font.name = "Calibri"
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = gtbl.cell(ri + 1, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = P_CARD if ri % 2 == 0 else P_BG
            c.margin_top = PInches(0.04); c.margin_bottom = PInches(0.04)
            c.margin_left = PInches(0.08); c.margin_right = PInches(0.08)
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = val
            r.font.size = PPt(fsize); r.font.color.rgb = P_TEXT; r.font.name = "Calibri"
    return gtbl


def _slide(prs, blank, kicker, n, title):
    s = prs.slides.add_slide(blank)
    _bg(s)
    _chrome(s, kicker, n)
    _h(s, title)
    return s


def build_pptx():
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    # ===================================================== TITELFOLIE (keine Nummer)
    s = prs.slides.add_slide(blank)
    _bg(s)
    _rect(s, 0, 3.05, 13.333, 0.06, P_CYAN)
    tf = _box(s, 0.9, 1.45, 11.5, 0.6)
    _set(tf, "ENGINEERING HARNESS", size=16, color=P_CYAN, bold=True)
    tf = _box(s, 0.9, 1.95, 11.5, 1.2)
    _set(tf, "ECC × BestPractice", size=46, color=P_WHITE, bold=True)
    tf = _box(s, 0.9, 3.25, 11.5, 1.6)
    _bullet(tf, "Das fusionierte, VPS-weite Claude-Code-Harness", size=20, color=P_TEXT, bold=True, first=True)
    _bullet(tf, "Global aktiv in ~/.claude · ECC 2.0.0-rc.1 (core, führend) · RTK-sicher", size=14, color=P_MUTED)
    _bullet(tf, "Quelle/Re-Install: ./install-vps.sh · Guides: @affaanmustafa", size=14, color=P_MUTED)

    # ===================================================== 1 · Überblick
    s = _slide(prs, blank, "Überblick", 1, "§1 Was ist das Harness?")
    tf = _box(s, 0.7, 1.9, 7.4, 4.5)
    _bullet(tf, "Fusion zweier Systeme — Fokus auf ECC:", size=17, color=P_INDIGO, bold=True, first=True)
    _bullet(tf, "ECC-Engine: 63 Agents · 249 Skills · ~80 Commands,", size=14, color=P_TEXT, level=1)
    _bullet(tf, "manifest-Installer, Continuous Learning, Security-Scanning.", size=14, color=P_MUTED, level=2)
    _bullet(tf, "BestPractice-Extras: /ecc-onboard (One-Shot-Setup),", size=14, color=P_TEXT, level=1)
    _bullet(tf, "state/-Pattern, Karpathy-Prinzipien.", size=14, color=P_MUTED, level=2)
    # Amber Zitat-Karte
    _round(s, 8.4, 1.9, 4.3, 2.6, P_CARD, line=P_AMBER)
    tf = _box(s, 8.7, 2.15, 3.7, 2.2)
    _bullet(tf, "„Konfiguration ist Fine-Tuning,", size=17, color=P_AMBER, bold=True, first=True)
    _bullet(tf, "nicht Architektur.“", size=17, color=P_AMBER, bold=True)
    _bullet(tf, "", size=6)
    _bullet(tf, "Muster wiederverwenden · Kontext sauber · billigstes Modell · Evidenz.", size=12, color=P_MUTED)
    tf = _box(s, 0.7, 6.4, 12.0, 0.7)
    _set(tf, "5 Prinzipien:  Agent-First · Test-Driven · Security-First · Immutability · Plan-Before-Execute.",
         size=14, color=P_INDIGO, bold=True)

    # ===================================================== 2 · Setup & VPS
    s = _slide(prs, blank, "Installation", 2, "§2 Setup & VPS-Architektur")
    tf = _box(s, 0.7, 1.85, 12.0, 1.7)
    _bullet(tf, "ECC ist GLOBAL (~/.claude) → aktiv in JEDEM VPS-Projekt (WMS, Jarvis, n8n …).",
            size=16, color=P_INDIGO, bold=True, first=True)
    _bullet(tf, "Namespace-sicher (rules/ecc, skills/ecc) · RTK-Hook unangetastet · Installer additiv.", size=14, color=P_TEXT)
    _bullet(tf, "Reproduzierbar: Repo = Quelle, Re-Install via ./install-vps.sh (idempotent, mit Backup).", size=14, color=P_TEXT)
    _dtable(s, ["Dein Tool", "ECC-Pendant", "Empfehlung"], [
        ["superpowers", "tdd-workflow, verification-loop", "beides nutzbar; ECC bevorzugen"],
        ["claude-mem", "/save-session, learning-v2", "primäre Memory-Quelle, auto-injiziert"],
        ["RTK", "Token-Ökonomie", "ergänzen sich (Bash- vs Workflow-Ebene)"],
        ["codex / superpowers", "diverse Skills", "eigenständig; bleiben aktiv"],
    ], 0.7, 3.7, 12.0, 2.8, col_w=[2.6, 4.2, 5.2], fsize=12)

    # ===================================================== 3 · Mentales Modell (Karten)
    s = _slide(prs, blank, "Bausteine", 3, "§3 Mentales Modell — die Bausteine")
    positions = [
        (0.7, 1.95), (3.85, 1.95), (7.0, 1.95), (10.15, 1.95),
        (0.7, 3.85), (3.85, 3.85), (7.0, 3.85),
    ]
    cw, ch = 2.95, 1.7
    for (bx, by), block in zip(positions, BUILDING_BLOCKS):
        _round(s, bx, by, cw, ch, P_CARD)
        tf = _box(s, bx + 0.18, by + 0.14, cw - 0.36, ch - 0.28)
        _bullet(tf, block[0], size=15, color=P_CYAN, bold=True, first=True)
        _bullet(tf, block[1], size=10.5, color=P_MUTED)
        _bullet(tf, block[2], size=9.5, color=P_INDIGO)
    tf = _box(s, 0.7, 5.9, 12.0, 0.9)
    _set(tf, "Regel:  Durable Logik → Skills.  3× getippt → /skill-create.",
         size=15, color=P_AMBER, bold=True)

    # ===================================================== 4 · Agent- & Skill-Katalog (NEU)
    s = _slide(prs, blank, "Katalog", 4, "§4 Agent- & Skill-Katalog")
    cards = [
        ("Core-Workflow-Agents", [a for a, _ in AGENTS_CORE]),
        ("Sprach-Reviewer", AGENTS_LANG),
        ("Build-Resolver", AGENTS_BUILD),
        ("Schlüssel-Skills", ["tdd-workflow", "verification-loop", "code-tour", "production-audit",
                              "eval-harness", "iterative-retrieval", "continuous-learning-v2",
                              "strategic-compact", "configure-ecc", "council"]),
    ]
    quad = [(0.7, 1.9), (7.0, 1.9), (0.7, 4.25), (7.0, 4.25)]
    cw, ch = 5.65, 2.15
    for (bx, by), (title, names) in zip(quad, cards):
        _round(s, bx, by, cw, ch, P_CARD)
        tf = _box(s, bx + 0.2, by + 0.15, cw - 0.4, ch - 0.3)
        _bullet(tf, title, size=14, color=P_CYAN, bold=True, first=True)
        _bullet(tf, " · ".join(names), size=10.5, color=P_TEXT)
    tf = _box(s, 0.7, 6.55, 12.0, 0.6)
    _set(tf, "63 Agents · 21 Skills (core; 249 via --profile full) · claude-mem + superpowers koexistieren.",
         size=12, color=P_MUTED)

    # ===================================================== 5 · Workflow 5 Phasen (Flow)
    s = _slide(prs, blank, "Pipeline", 5, "§5 Workflow — die 5 Phasen")
    phases = [
        ("1", "RESEARCH", "Explore · claude-mem", "→ research.md"),
        ("2", "PLAN", "/plan", "→ plan.md"),
        ("3", "IMPLEMENT", "/feature-dev + tdd-workflow", "→ Code + Tests"),
        ("4", "REVIEW", "/code-review", "→ review.md"),
        ("5", "VERIFY", "Tests · Build · /build-fix", "→ grün ↺"),
    ]
    node_w, node_h = 2.18, 2.3
    gap = 0.32
    x = 0.55
    y = 2.2
    for i, (num, name, tool, out) in enumerate(phases):
        _round(s, x, y, node_w, node_h, P_CARD, line=P_CYAN)
        tf = _box(s, x + 0.12, y + 0.18, node_w - 0.24, node_h - 0.3)
        _bullet(tf, num, size=30, color=P_CYAN, bold=True, first=True)
        _bullet(tf, name, size=14, color=P_TEXT, bold=True)
        _bullet(tf, tool, size=10, color=P_MUTED)
        _bullet(tf, out, size=10, color=P_CYAN)
        # Chevron-Verbinder zwischen Knoten
        if i < len(phases) - 1:
            ch_x = x + node_w + (gap - 0.28) / 2
            chev = s.shapes.add_shape(MSO_SHAPE.CHEVRON, PInches(ch_x), PInches(y + node_h / 2 - 0.18),
                                      PInches(0.28), PInches(0.36))
            chev.fill.solid(); chev.fill.fore_color.rgb = P_CYAN
            chev.line.fill.background(); chev.shadow.inherit = False
        x += node_w + gap
    tf = _box(s, 0.7, 4.95, 12.0, 2.0)
    _bullet(tf, "Ein Input, ein Output je Schritt — in Dateien speichern.", size=14, color=P_MUTED, first=True)
    _bullet(tf, "/clear zwischen großen Phasen — Explorations-Kontext ist für die Umsetzung irrelevant.", size=14, color=P_MUTED)
    _bullet(tf, "Subagents immer objektiven Kontext mitgeben, nicht nur die Frage.", size=14, color=P_MUTED)

    # ===================================================== 6 · /ecc-onboard
    s = _slide(prs, blank, "Onboarding", 6, "§6a Projekt-Onboarding mit /ecc-onboard")
    steps = [
        ("1 · Detect", "Stack erkennen (package.json, go.mod …)", "Profil + Sprach-Packs"),
        ("2 · Dry-Run", "install-apply --dry-run", "Plan ohne Schreiben"),
        ("3 · OK", "1× Bestätigung des Users", "Freigabe"),
        ("4 · Install", "install-apply --target claude-project", ".claude/rules/ecc + skills/ecc"),
        ("5 · Kontext", "Templates mit echten Werten füllen", "PROJECT_RULES.md + state/"),
    ]
    cw = 2.34
    x = 0.55
    y = 2.15
    for title, what, res in steps:
        _round(s, x, y, cw, 3.0, P_CARD, line=P_CYAN)
        tf = _box(s, x + 0.16, y + 0.18, cw - 0.32, 2.7)
        _bullet(tf, title, size=14, color=P_CYAN, bold=True, first=True)
        _bullet(tf, what, size=11, color=P_TEXT)
        _bullet(tf, "", size=4)
        _bullet(tf, "→ " + res, size=11, color=P_INDIGO)
        x += cw + 0.18
    tf = _box(s, 0.7, 5.6, 12.0, 1.2)
    _bullet(tf, "Aus dem Zielprojekt ausführen — Install landet im aktuellen cwd.", size=15, color=P_AMBER, bold=True, first=True)
    _bullet(tf, "Idempotent: erneut ausführen merged, überschreibt nichts blind.", size=14, color=P_MUTED)

    # ===================================================== 7 · Codemaps (KERN-FEATURE)
    s = _slide(prs, blank, "Codemaps", 7, "§6b Codebase EINMAL erkunden")
    tf = _box(s, 0.7, 1.85, 12.0, 0.9)
    _set(tf, "Prinzip: einmal kartieren → in Dateien festhalten → künftig die Karte lesen statt neu suchen.",
         size=15, color=P_INDIGO, bold=True)
    _dtable(s, ["Datei (docs/CODEMAPS/)", "Inhalt"], [
        ["architecture.md", "System-Diagramm, Service-Grenzen, Datenfluss"],
        ["backend.md", "Routes → Controller → Service → Repository"],
        ["frontend.md", "Page-Tree, Komponenten, State-Flow"],
        ["data.md", "Tabellen, Relationen, Migrationen"],
        ["dependencies.md", "externe Dienste, Integrationen"],
    ], 0.7, 2.75, 12.0, 2.9, col_w=[3.3, 8.7], fsize=12)
    tf = _box(s, 0.7, 5.85, 12.0, 1.3)
    _bullet(tf, "Snapshot — nach Architektur-Änderung neu bauen. Fremdes Repo: code-tour / claude-mem smart-explore.",
            size=14, color=P_AMBER, bold=True, first=True)
    _bullet(tf, "Pro Aufgabe: research.md (Explore-Agent) → /plan. claude-mem injiziert Alt-Kontext automatisch.",
            size=13, color=P_MUTED)

    # ===================================================== 8 · §7a Zwei Startpunkte (NEU)
    s = _slide(prs, blank, "Start", 8, "§7a Zwei Startpunkte — leer oder Codebase")
    _round(s, 0.7, 1.95, 5.85, 3.5, P_CODE)
    tf = _box(s, 0.95, 2.18, 5.4, 3.2)
    _bullet(tf, "LEERES PROJEKT (greenfield)", size=14, color=P_CYAN, bold=True, first=True)
    _bullet(tf, "", size=6)
    for line in [
            "mkdir app && cd app && git init",
            "claude",
            "/ecc-onboard      # Stack → 1x OK",
            "/plan  Baue <Feature> …"]:
        _bullet(tf, line, size=12.5, color=PRGB(0x9C, 0xE3, 0xF0), font="Consolas")
    _round(s, 6.85, 1.95, 5.8, 3.5, P_CODE)
    tf = _box(s, 7.1, 2.18, 5.35, 3.2)
    _bullet(tf, "BESTEHENDE CODEBASE", size=14, color=P_CYAN, bold=True, first=True)
    _bullet(tf, "", size=6)
    for line in [
            "git clone <repo> && cd <repo>",
            "claude",
            "/ecc-onboard      # Stack erkannt",
            "/update-codemaps  # EINMAL kartieren",
            "code-tour · smart-explore"]:
        _bullet(tf, line, size=12.5, color=PRGB(0x9C, 0xE3, 0xF0), font="Consolas")
    tf = _box(s, 0.7, 5.7, 12.0, 1.3)
    _bullet(tf, "Einmal, nicht jedes Mal: /ecc-onboard + /update-codemaps laufen pro Projekt EINMAL.",
            size=15, color=P_AMBER, bold=True, first=True)
    _bullet(tf, "Ergebnis ist persistent (Dateien) — jede Folge-Session startet mit der kompakten Karte.",
            size=13, color=P_MUTED)

    # ===================================================== 9 · §7b Feature durchgespielt (NEU)
    s = _slide(prs, blank, "Durchgespielt", 9, "§7b Ein Feature von 0 auf grün — /api/search")
    _dtable(s, ["Phase", "Du tippst", "→ Output · Modell"], [
        ["1 RESEARCH", "Explore-Agent: Routing/Handler/Tests kartieren (+ claude-mem)", "research.md · Haiku"],
        ["2 PLAN", "/plan  /api/search: Query → Filter → Paginierung, Tests zuerst", "plan.md (wartet) · Opus"],
        ["3 IMPLEMENT", "/feature-dev  plan.md umsetzen (tdd-workflow: Tests zuerst)", "Code+Tests RED→GREEN · Sonnet"],
        ["4 REVIEW", "/code-review  (+ /python-review …)", "review.md · Opus"],
        ["5 VERIFY", "Tests/Build grün? sonst /build-fix → zurück zu Phase 3", "grün → commit · Sonnet"],
    ], 0.6, 1.95, 12.2, 3.6, col_w=[2.0, 6.1, 4.1], fsize=11.5)
    tf = _box(s, 0.7, 5.75, 12.0, 1.3)
    _bullet(tf, "/clear zwischen den Phasen. Jede Phase: ein Input, ein Output (Datei). Phasen nicht überspringen.",
            size=14, color=P_AMBER, bold=True, first=True)
    _bullet(tf, "Regel: „In Dateien speichern, nicht im Kopf behalten.“  research.md → plan.md → review.md.",
            size=13, color=P_MUTED)

    # ===================================================== 10 · Cheat-Sheet 1/2
    s = _slide(prs, blank, "Cheat-Sheet", 10, "§8 Cheat-Sheet — Command (1/2)")
    _dtable(s, ["Situation", "Command", "Zweck"], CMD_CHEAT[:8], 0.6, 1.9, 12.2, 5.0,
            col_w=[3.0, 4.4, 4.8], fsize=11.5)

    # ===================================================== 9 · Cheat-Sheet 2/2
    s = _slide(prs, blank, "Cheat-Sheet", 11, "§8 Cheat-Sheet — Command (2/2)")
    _dtable(s, ["Situation", "Command", "Zweck"], CMD_CHEAT[8:], 0.6, 1.9, 12.2, 5.0,
            col_w=[3.0, 4.4, 4.8], fsize=11.5)

    # ===================================================== 10 · MCP
    s = _slide(prs, blank, "MCP", 12, "§9 Wann welcher MCP")
    tf = _box(s, 0.7, 1.8, 12.0, 0.45)
    _set(tf, "Faustregel: < 10 MCPs aktiv / < 80 Tools. Ungenutztes mit /mcp deaktivieren.",
         size=14, color=P_AMBER, bold=True)
    tf = _box(s, 0.7, 2.3, 12.0, 0.38)
    _set(tf, "● BEI DIR AKTIV", size=13, color=P_CYAN, bold=True)
    _dtable(s, ["Aufgabe", "MCP / Tool", "Hinweis"], MCP_ACTIVE, 0.6, 2.72, 12.2, 3.35,
            col_w=[3.4, 2.9, 5.9], fsize=11)
    tf = _box(s, 0.7, 6.2, 12.0, 1.05)
    _bullet(tf, "Optional, NICHT installiert: exa · github (du: gh-CLI) · playwright · sequential-thinking.",
            size=13, color=P_MUTED, first=True)
    _bullet(tf, "„memory“-MCP (Wissensgraph) nicht nötig — claude-mem ist deine Memory-Quelle.",
            size=13, color=P_AMBER, bold=True)

    # ===================================================== 11 · Modellstrategie (Karten)
    s = _slide(prs, blank, "Modelle", 13, "§10 Token-Ökonomie & Modellstrategie")
    mcards = [
        ("HAIKU 4.5", P_CYAN, "Worker-Agents", [
            "Exploration / Suche", "viele günstige Parallel-Agents", "Doku-Struktur"]),
        ("SONNET 4.6", P_INDIGO, "günstige Alternative", [
            "einfache, gut definierte Edits", "PR-Review (Kontext)", "Doku schreiben"]),
        ("OPUS 4.8 (1M)", P_AMBER, "Standard / Default", [
            "Standard-Coding / Multi-File (/fast)", "komplexe Architektur (effort high)", "Security-Analyse · komplexe Bugs"]),
    ]
    x = 0.7
    cw = 3.9
    for name, accent, role, pts in mcards:
        _round(s, x, 1.95, cw, 3.7, P_CARD, line=accent)
        tf = _box(s, x + 0.22, 2.15, cw - 0.44, 3.4)
        _bullet(tf, name, size=18, color=accent, bold=True, first=True)
        _bullet(tf, role, size=13, color=P_TEXT, bold=True)
        _bullet(tf, "", size=6)
        for pt in pts:
            _bullet(tf, "• " + pt, size=12, color=P_MUTED)
        x += cw + 0.22
    tf = _box(s, 0.7, 5.95, 12.0, 1.0)
    _set(tf, "Effort medium Standard / high bei Architektur+Security · /model-route · /cost-report · Auto-Compact aus.",
         size=13, color=P_MUTED)

    # ===================================================== 12 · Memory + Learning
    s = _slide(prs, blank, "Memory & Lernen", 14, "§11–12 Memory, Sessions & Continuous Learning")
    _dtable(s, ["Command", "Zweck"], [
        ["/save-session · /resume-session", "Tag-zu-Tag-Fortschritt persistieren"],
        ["/sessions · /checkpoint · /aside", "Historie · Checkpoint · Nebenfrage"],
        ["/learn-eval → /evolve → /promote", "Self-Improvement-Loop (Instincts)"],
        ["/instinct-status · /skill-create", "Instincts zeigen · Skill aus Git-History"],
    ], 0.7, 1.95, 12.0, 3.0, col_w=[5.4, 6.6], fsize=13)
    tf = _box(s, 0.7, 5.2, 12.0, 1.7)
    _bullet(tf, "claude-mem = deine Memory-Quelle: Auto-Injektion ab 2. Session · /mem-search für frühere Arbeit.",
            size=14, color=P_TEXT, bold=True, first=True)
    _bullet(tf, "EINE primäre Memory-Quelle je Projekt — kein zusätzlicher memory-MCP. Keine Secrets in Memory.",
            size=14, color=P_AMBER, bold=True)
    _bullet(tf, "Routine: Session-Ende → /learn-eval → /evolve → /promote.", size=13, color=P_INDIGO)

    # ===================================================== 13 · Parallelisierung (NEU)
    s = _slide(prs, blank, "Parallel", 15, "§13 Parallelisierung")
    pcards = [
        ("/fork", "Konversation forken für nicht-überlappende Nebenaufgaben (Codebase-Fragen)."),
        ("Git-Worktrees", "überlappende parallele Arbeit ohne Konflikte — eigene Claude-Instanz je Worktree."),
        ("Cascade-Methode", "neue Tasks in neuen Tabs, alt→neu abarbeiten, max. 3–4 gleichzeitig."),
        ("ECC-Multi-Model", "/multi-plan · /multi-workflow · /multi-backend · /multi-frontend · /loop-start."),
    ]
    quad = [(0.7, 1.95), (7.0, 1.95), (0.7, 3.65), (7.0, 3.65)]
    cw, ch = 5.65, 1.5
    for (bx, by), (title, desc) in zip(quad, pcards):
        _round(s, bx, by, cw, ch, P_CARD, line=P_CYAN)
        tf = _box(s, bx + 0.2, by + 0.15, cw - 0.4, ch - 0.3)
        _bullet(tf, title, size=15, color=P_CYAN, bold=True, first=True)
        _bullet(tf, desc, size=11.5, color=P_MUTED)
    _round(s, 0.7, 5.4, 11.95, 1.05, P_CODE)
    tf = _box(s, 0.95, 5.55, 11.5, 0.8)
    _bullet(tf, "git worktree add ../feature-a feature-a && cd ../feature-a && claude",
            size=14, color=PRGB(0x9C, 0xE3, 0xF0), bold=True, first=True, font="Consolas")

    # ===================================================== 14 · Security
    s = _slide(prs, blank, "Security", 16, "§14 Security im Alltag + VPS")
    _round(s, 0.7, 1.95, 5.85, 3.7, P_CARD, line=P_EMERALD)
    tf = _box(s, 0.95, 2.15, 5.4, 3.4)
    _bullet(tf, "Aktiv / sicher", size=16, color=P_EMERALD, bold=True, first=True)
    _bullet(tf, "", size=4)
    for lab, txt in SECURITY_ACTIVE:
        _bullet(tf, lab, size=12.5, color=P_TEXT, bold=True)
        _bullet(tf, txt, size=11, color=P_MUTED, level=1)
    _round(s, 6.85, 1.95, 5.8, 3.7, P_CARD, line=P_AMBER)
    tf = _box(s, 7.1, 2.15, 5.35, 3.4)
    _bullet(tf, "Bewusst opt-in", size=16, color=P_AMBER, bold=True, first=True)
    _bullet(tf, "", size=4)
    for lab, txt in SECURITY_OPTIN:
        _bullet(tf, lab, size=12.5, color=P_TEXT, bold=True)
        _bullet(tf, txt, size=11, color=P_MUTED, level=1)
    tf = _box(s, 0.7, 5.85, 12.0, 1.1)
    _set(tf, "VPS-weit: ECC global → Hooks bleiben inaktiv bis zur bewussten Aktivierung. "
             "Härtung: ./install-vps.sh --harden (mit settings.json-Backup).", size=13, color=P_MUTED)

    # ===================================================== 15 · ECC-CLIs & Wartung (NEU)
    s = _slide(prs, blank, "CLIs & Wartung", 17, "§15 ECC-CLIs & Wartung")
    _round(s, 0.7, 1.95, 11.95, 3.2, P_CODE)
    tf = _box(s, 1.0, 2.2, 11.4, 2.8)
    code_lines = [
        "node scripts/ecc.js doctor          # Health-Check (claude-home: OK)",
        "node scripts/ecc.js list-installed   # was ist installiert",
        "node scripts/ecc.js repair           # ECC-Dateien wiederherstellen",
        'npx ecc consult "security reviews"   # Advisor: welche Komponenten?',
        "./install-vps.sh                     # Re-Install / Update VPS-weit",
    ]
    for i, line in enumerate(code_lines):
        _bullet(tf, line, size=13, color=PRGB(0x9C, 0xE3, 0xF0), first=(i == 0), font="Consolas")
    tf = _box(s, 0.7, 5.5, 12.0, 1.2)
    _set(tf, "Slash-Einstieg:  /ecc-guide · /harness-audit · /cost-report · /skill-health · /projects.",
         size=15, color=P_INDIGO, bold=True)

    # ===================================================== 16 · Quellen & TL;DR
    s = _slide(prs, blank, "Quellen", 18, "Quellen & TL;DR")
    tf = _box(s, 0.7, 1.9, 12.0, 1.7)
    _bullet(tf, "Guides von @affaanmustafa (X-Threads; lokal unter ecc/):", size=14, color=P_CYAN, bold=True, first=True)
    _bullet(tf, "Shortform: " + X_SHORT, size=11.5, color=P_TEXT)
    _bullet(tf, "Longform:  " + X_LONG, size=11.5, color=P_TEXT)
    _bullet(tf, "Security:  " + X_SEC, size=11.5, color=P_TEXT)
    tf = _box(s, 0.7, 3.7, 12.0, 3.4)
    _bullet(tf, "TL;DR — wie der Erfinder:", size=16, color=P_CYAN, bold=True, first=True)
    _bullet(tf, "1. In Phasen: Research → /plan → /feature-dev+TDD → /code-review → Verify.", size=13, color=P_TEXT)
    _bullet(tf, "2. Kontext sauber: <10 MCPs, billigstes Modell, modulare Dateien, strategisch compacten.", size=13, color=P_TEXT)
    _bullet(tf, "3. Persistieren: /save-session · /resume-session.", size=13, color=P_TEXT)
    _bullet(tf, "4. Lernen: /learn-eval → /evolve → /promote.", size=13, color=P_TEXT)
    _bullet(tf, "5. Security zuerst: RTK aktiv, Hooks/AgentShield bewusst, nie blind ausführen.", size=13, color=P_TEXT)
    _bullet(tf, "6. Unsicher? /ecc-guide · ecc doctor · die drei Original-Guides.", size=13, color=P_TEXT)

    prs.save(PPTX_OUT)
    print("✅ PowerPoint:", PPTX_OUT)


if __name__ == "__main__":
    build_docx()
    build_pptx()
