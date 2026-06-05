#!/usr/bin/env python3
"""Build-Pipeline: docs/ECC-Harness-Guide.de.md -> .docx + .pptx.

Schicht-2-Tool. Single Source ist die Markdown-Datei; docx/pptx werden daraus
reproduzierbar erzeugt, damit die Office-Snapshots nicht mehr von Hand driften.

Unterstütztes Markdown-Subset (bewusst klein gehalten):
  # Titel              -> Titelseite (genau einmal, ganz oben)
  > ... (vor 1. ##)    -> Untertitel-Block der Titelseite
  ## N · Kapitel       -> Kapitel  (docx Heading 1 / je eine pptx-Folie)
  ### Unterabschnitt   -> docx Heading 2 / fette Zeile auf der Folie
  | a | b |            -> GFM-Tabelle (Header + ----Trennzeile + Zeilen)
  - item               -> Aufzählung
  ```                  -> Codeblock (monospace)
  > callout            -> Hinweis-Box (innerhalb eines Kapitels)
  **bold** *italic* `code`  -> inline

Aufruf:
  python3 build.py                # baut beide Snapshots neben die md-Quelle
  python3 build.py --check        # nur Parsen + Statistik, nichts schreiben
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

# --- Pfade -----------------------------------------------------------------
REPO_ROOT = Path(__file__).resolve().parents[3]
DOCS = REPO_ROOT / "docs"
SRC = DOCS / "ECC-Harness-Guide.de.md"
OUT_DOCX = DOCS / "ECC-Harness-Guide.de.docx"
OUT_PPTX = DOCS / "ECC-Harness-Guide.de.pptx"

# --- Farben / Stil ---------------------------------------------------------
ACCENT = (0x1F, 0x6F, 0x8B)        # Teal-Akzent
DARK = (0x0F, 0x2A, 0x3C)          # Titelfolie-Hintergrund
LIGHT = (0xF4, 0xF6, 0xF8)         # Codeblock-Schattierung (docx fill hex)
MONO_FONT = "Consolas"
BODY_FONT = "Calibri"


# ==========================================================================
# 1) Markdown-Parser  ->  Liste von Block-dicts
# ==========================================================================
INLINE_RE = re.compile(r"(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`)")


def parse_inline(text: str):
    """Zerlegt einen String in Runs: (text, bold, italic, mono)."""
    runs = []
    for tok in INLINE_RE.split(text):
        if not tok:
            continue
        if tok.startswith("**") and tok.endswith("**"):
            runs.append((tok[2:-2], True, False, False))
        elif tok.startswith("*") and tok.endswith("*"):
            runs.append((tok[1:-1], False, True, False))
        elif tok.startswith("`") and tok.endswith("`"):
            runs.append((tok[1:-1], False, False, True))
        else:
            runs.append((tok, False, False, False))
    return runs


def split_table_row(line: str):
    cells = [c.strip() for c in line.strip().strip("|").split("|")]
    return cells


def parse(md: str):
    lines = md.split("\n")
    blocks = []
    i = 0
    n = len(lines)
    while i < n:
        line = lines[i]
        stripped = line.strip()

        # Leerzeile
        if not stripped:
            i += 1
            continue

        # Horizontale Linie
        if stripped == "---":
            blocks.append({"type": "hr"})
            i += 1
            continue

        # Codeblock
        if stripped.startswith("```"):
            code = []
            i += 1
            while i < n and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1  # schließendes ```
            blocks.append({"type": "code", "lines": code})
            continue

        # Überschriften
        if stripped.startswith("### "):
            blocks.append({"type": "h2", "text": stripped[4:].strip()})
            i += 1
            continue
        if stripped.startswith("## "):
            blocks.append({"type": "h1", "text": stripped[3:].strip()})
            i += 1
            continue
        if stripped.startswith("# "):
            blocks.append({"type": "title", "text": stripped[2:].strip()})
            i += 1
            continue

        # Blockquote (mehrzeilig)
        if stripped.startswith(">"):
            quote = []
            while i < n and lines[i].strip().startswith(">"):
                q = lines[i].strip()[1:].lstrip()
                quote.append(q)
                i += 1
            blocks.append({"type": "quote", "lines": [q for q in quote if q != ""]})
            continue

        # Tabelle (Header | Trennzeile | Zeilen)
        if stripped.startswith("|") and i + 1 < n and set(lines[i + 1].strip()) <= set("|-: "):
            header = split_table_row(lines[i])
            i += 2  # Header + Trennzeile
            rows = []
            while i < n and lines[i].strip().startswith("|"):
                rows.append(split_table_row(lines[i]))
                i += 1
            blocks.append({"type": "table", "header": header, "rows": rows})
            continue

        # Liste
        if stripped.startswith("- ") or stripped.startswith("* "):
            items = []
            while i < n and (lines[i].strip().startswith("- ") or lines[i].strip().startswith("* ")):
                items.append(lines[i].strip()[2:].strip())
                i += 1
            blocks.append({"type": "bullets", "items": items})
            continue

        # Absatz (bis Leerzeile / Strukturzeile)
        para = [stripped]
        i += 1
        while i < n:
            s = lines[i].strip()
            if (not s or s.startswith("#") or s.startswith("|") or s.startswith("- ")
                    or s.startswith("* ") or s.startswith("```") or s.startswith(">")
                    or s == "---"):
                break
            para.append(s)
            i += 1
        blocks.append({"type": "para", "text": " ".join(para)})
    return blocks


# ==========================================================================
# 2) DOCX-Renderer
# ==========================================================================
def build_docx(blocks):
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = BODY_FONT
    normal.font.size = Pt(10.5)

    def fill_runs(paragraph, text, base_size=10.5):
        for txt, bold, italic, mono in parse_inline(text):
            run = paragraph.add_run(txt)
            run.bold = bold
            run.italic = italic
            if mono:
                run.font.name = MONO_FONT
                run.font.size = Pt(base_size - 1)

    def shade(paragraph, fill_hex):
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:fill"), fill_hex)
        paragraph.paragraph_format.element.get_or_add_pPr().append(shd)

    # --- Titelseite ---
    title_block = next((b for b in blocks if b["type"] == "title"), None)
    first_h1 = next((idx for idx, b in enumerate(blocks) if b["type"] == "h1"), len(blocks))

    if title_block:
        t = doc.add_paragraph()
        t.alignment = WD_ALIGN_PARAGRAPH.LEFT
        r = t.add_run(title_block["text"])
        r.bold = True
        r.font.size = Pt(30)
        r.font.color.rgb = RGBColor(*ACCENT)
        t.paragraph_format.space_before = Pt(120)

    # Untertitel = quote-Blöcke vor dem ersten Kapitel
    for b in blocks[:first_h1]:
        if b["type"] == "quote":
            for ln in b["lines"]:
                p = doc.add_paragraph()
                fill_runs(p, ln, base_size=11)
                p.paragraph_format.space_after = Pt(4)

    doc.add_page_break()

    # --- Inhaltsverzeichnis (echtes TOC-Feld) ---
    h = doc.add_heading("Inhalt", level=1)
    toc_p = doc.add_paragraph()
    run = toc_p.add_run()
    for el, attr, val, txt in (
        ("w:fldChar", "w:fldCharType", "begin", None),
        ("w:instrText", "xml:space", "preserve", 'TOC \\o "1-2" \\h \\z \\u'),
        ("w:fldChar", "w:fldCharType", "separate", None),
        ("w:t", None, None, "Inhaltsverzeichnis – in Word mit F9 aktualisieren."),
        ("w:fldChar", "w:fldCharType", "end", None),
    ):
        e = OxmlElement(el)
        if attr:
            e.set(qn(attr), val)
        if txt is not None:
            e.text = txt
        run._r.append(e)
    doc.add_page_break()

    # --- Kapitel ---
    for b in blocks[first_h1:]:
        bt = b["type"]
        if bt == "h1":
            doc.add_heading(b["text"], level=1)
        elif bt == "h2":
            doc.add_heading(b["text"], level=2)
        elif bt == "para":
            p = doc.add_paragraph()
            fill_runs(p, b["text"])
        elif bt == "bullets":
            for it in b["items"]:
                p = doc.add_paragraph(style="List Bullet")
                fill_runs(p, it)
        elif bt == "code":
            for ln in b["lines"]:
                p = doc.add_paragraph()
                shade(p, LIGHT[0] if isinstance(LIGHT, str) else "F4F6F8")
                r = p.add_run(ln if ln else " ")
                r.font.name = MONO_FONT
                r.font.size = Pt(9)
        elif bt == "quote":
            for ln in b["lines"]:
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Inches(0.25)
                shade(p, "EAF2F5")
                fill_runs(p, ln)
        elif bt == "table":
            cols = max(len(b["header"]), max((len(r) for r in b["rows"]), default=0))
            table = doc.add_table(rows=1, cols=cols)
            table.style = "Light Grid Accent 1"
            hdr = table.rows[0].cells
            for j in range(cols):
                cell = hdr[j]
                cell.paragraphs[0].text = ""
                rr = cell.paragraphs[0].add_run(b["header"][j] if j < len(b["header"]) else "")
                rr.bold = True
            for row in b["rows"]:
                cells = table.add_row().cells
                for j in range(cols):
                    cells[j].paragraphs[0].text = ""
                    fill_runs(cells[j].paragraphs[0], row[j] if j < len(row) else "")
            doc.add_paragraph()

    doc.save(str(OUT_DOCX))
    return OUT_DOCX


# ==========================================================================
# 3) PPTX-Renderer (eine Folie je Kapitel, Überlauf -> Folgefolie)
# ==========================================================================
def build_pptx(blocks):
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    def add_rect(slide, x, y, w, h, rgb):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(*rgb)
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    # --- Titelfolie ---
    title_block = next((b for b in blocks if b["type"] == "title"), None)
    first_h1 = next((idx for idx, b in enumerate(blocks) if b["type"] == "h1"), len(blocks))
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, DARK)
    add_rect(s, Inches(0), Inches(3.1), SW, Inches(0.08), ACCENT)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_block["text"] if title_block else "Engineering Harness"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Untertitelzeilen
    sub = s.shapes.add_textbox(Inches(0.9), Inches(3.4), Inches(11.5), Inches(3.4))
    stf = sub.text_frame
    stf.word_wrap = True
    first = True
    for b in blocks[:first_h1]:
        if b["type"] == "quote":
            for ln in b["lines"]:
                para = stf.paragraphs[0] if first else stf.add_paragraph()
                first = False
                plain = re.sub(r"[*`]", "", ln)
                para.text = plain
                para.font.size = Pt(13)
                para.font.color.rgb = RGBColor(0xD7, 0xE3, 0xEA)
                para.space_after = Pt(6)

    def new_content_slide(title, cont=False):
        sl = prs.slides.add_slide(BLANK)
        add_rect(sl, 0, 0, SW, Inches(0.95), ACCENT)
        tbox = sl.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.7))
        ttf = tbox.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tr = tp.add_run()
        tr.text = title + (" (Forts.)" if cont else "")
        tr.font.size = Pt(24)
        tr.font.bold = True
        tr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        body = sl.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12.1), Inches(6.0))
        btf = body.text_frame
        btf.word_wrap = True
        return btf

    def flush_items(title, items):
        """items: Liste von (text, level, mono, bold). Splittet bei Überlauf."""
        if not items:
            items = [("—", 0, False, False)]
        # dynamische Schriftgröße + Folien-Kapazität
        per_slide = 13
        chunks = [items[k:k + per_slide] for k in range(0, len(items), per_slide)]
        for ci, chunk in enumerate(chunks):
            btf = new_content_slide(title, cont=(ci > 0))
            size = 16 if len(chunk) <= 8 else (14 if len(chunk) <= 11 else 12)
            first_para = True
            for (text, level, mono, bold) in chunk:
                para = btf.paragraphs[0] if first_para else btf.add_paragraph()
                first_para = False
                para.level = min(level, 4)
                run = para.add_run()
                run.text = text
                run.font.size = Pt(size if level == 0 else size - 1)
                run.font.bold = bold
                if mono:
                    run.font.name = MONO_FONT
                    run.font.size = Pt(size - 2)
                para.space_after = Pt(3)

    # Kapitel sammeln
    cur_title = None
    items = []
    for b in blocks[first_h1:]:
        bt = b["type"]
        if bt == "h1":
            if cur_title is not None:
                flush_items(cur_title, items)
            cur_title = b["text"]
            items = []
        elif bt == "h2":
            items.append((b["text"], 0, False, True))
        elif bt == "para":
            items.append((re.sub(r"[*`]", "", b["text"]), 0, False, False))
        elif bt == "bullets":
            for it in b["items"]:
                items.append(("• " + re.sub(r"[*`]", "", it), 0, False, False))
        elif bt == "code":
            for ln in b["lines"]:
                if ln.strip():
                    items.append((ln, 1, True, False))
        elif bt == "quote":
            for ln in b["lines"]:
                items.append(("▸ " + re.sub(r"[*`]", "", ln), 0, False, False))
        elif bt == "table":
            items.append((" · ".join(b["header"]), 0, False, True))
            for row in b["rows"]:
                items.append((" — ".join(c for c in row if c), 1, False, False))
    if cur_title is not None:
        flush_items(cur_title, items)

    prs.save(str(OUT_PPTX))
    return OUT_PPTX, len(prs.slides._sldIdLst)


# ==========================================================================
def main():
    ap = argparse.ArgumentParser(description="Baut docx+pptx aus der md-Quelle.")
    ap.add_argument("--check", action="store_true", help="nur parsen, nichts schreiben")
    args = ap.parse_args()

    if not SRC.exists():
        sys.exit(f"FEHLER: Quelle nicht gefunden: {SRC}")

    blocks = parse(SRC.read_text(encoding="utf-8"))
    n_ch = sum(1 for b in blocks if b["type"] == "h1")
    n_tbl = sum(1 for b in blocks if b["type"] == "table")
    print(f"Quelle: {SRC.relative_to(REPO_ROOT)}  ({len(blocks)} Blöcke, {n_ch} Kapitel, {n_tbl} Tabellen)")

    if args.check:
        print("--check: nur Parsen, kein Schreiben.")
        return

    docx_path = build_docx(blocks)
    pptx_path, n_slides = build_pptx(blocks)
    print(f"✓ docx: {docx_path.relative_to(REPO_ROOT)}")
    print(f"✓ pptx: {pptx_path.relative_to(REPO_ROOT)}  ({n_slides} Folien)")


if __name__ == "__main__":
    main()
